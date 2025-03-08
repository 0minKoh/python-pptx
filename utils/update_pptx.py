from pptx import Presentation
from pptx.presentation import Presentation as PresentationType
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE


def load_template(template_path: str) -> PresentationType:
    return Presentation(template_path)

def edit_text_field(prs: PresentationType, slide_index: int, is_title: bool, new_text: str, align_center: bool = True, ph_index: int = 11) -> PresentationType:
    slide = prs.slides[slide_index]
    if is_title:
        title = slide.shapes.title
        title.text = new_text
        title.text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
        if align_center == False:
            title.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
    else:
        content = slide.shapes.placeholders[ph_index]
        content.text = new_text
    return prs

def _make_lyrics_data(splited_lyrics: list) -> list:

    pages = []

    for lyrics in splited_lyrics:
        lyrics_data = {
            "lyrics": ""
        }
        lyrics_data["lyrics"] = lyrics
        pages.append(lyrics_data)
        
    return pages


def duplicate_and_add_slide(prs: PresentationType, duplicate_slide_index: int, slide_texts: list) -> dict:
    lyrics_data = _make_lyrics_data(slide_texts)
    original_slide = prs.slides[duplicate_slide_index]

    i = 0
    is_first_lyrics = True

    for lyrics in lyrics_data:
        slide_text = lyrics["lyrics"]
        # 첫 번째 가사는 기존 페이지의 글자를 수정함
        if is_first_lyrics:
            prs = edit_text_field(
                prs=prs,
                slide_index=duplicate_slide_index,
                is_title=True,
                new_text=slide_text
            )
            is_first_lyrics = False
            continue

        # 새 슬라이드를 duplicate_slide_index 다음에 추가
        slide_layout = original_slide.slide_layout
        new_slide = prs.slides.add_slide(slide_layout)
        new_slide_index = duplicate_slide_index + i + 1

        old_index = prs.slides.index(new_slide)

        xml_slides = prs.slides._sldIdLst  # pylint: disable=W0212
        slides = list(xml_slides)
        xml_slides.remove(slides[old_index])
        xml_slides.insert(new_slide_index, slides[old_index])

        # 텍스트 필드 편집
        prs = edit_text_field(
            prs=prs,
            slide_index=new_slide_index,
            is_title = True,
            new_text=slide_text
        )
        i += 1

    return {
        "prs": prs,
        "added_slide_count": i
    }

def add_ads_slides(prs: PresentationType, ads: list, ad_slide_index: int) -> PresentationType:
    is_first_ad = True

    for index, ad in enumerate(ads):
        title = ad["title"]
        contents = ad["contents"]

        if is_first_ad:
            prs = edit_text_field(
                prs=prs,
                slide_index=ad_slide_index,
                is_title=True,
                new_text=title,
                align_center=False
            )
            prs = edit_text_field(
                prs=prs,
                slide_index=ad_slide_index,
                is_title=False,
                new_text=contents
            )
            is_first_ad = False
            continue

        slide_layout = prs.slides[ad_slide_index].slide_layout
        new_slide = prs.slides.add_slide(slide_layout)
        new_slide_index = ad_slide_index + index

        old_index = prs.slides.index(new_slide)
        xml_slides = prs.slides._sldIdLst  # pylint: disable=W0212
        slides = list(xml_slides)
        xml_slides.remove(slides[old_index])
        xml_slides.insert(new_slide_index, slides[old_index])

        # 텍스트 필드 편집
        prs = edit_text_field(
            prs=prs,
            slide_index=new_slide_index,
            is_title=True,
            new_text=title,
            align_center=False
        )
        prs = edit_text_field(
            prs=prs,
            slide_index=new_slide_index,
            is_title=False,
            new_text=contents
        )

    return prs

def add_bible_slides(prs: PresentationType, bible_data: list, bible_slide_index: int) -> PresentationType:
    is_first_bible = True

    for index, bible in enumerate(bible_data):
        title = bible["title"]
        contents = bible["contents"]
        if is_first_bible:
            prs = edit_text_field(
                prs=prs,
                slide_index=bible_slide_index,
                is_title=True,
                new_text=title,
            )
            prs = edit_text_field(
                prs=prs,
                slide_index=bible_slide_index,
                is_title=False,
                new_text=contents,
                ph_index=10
            )
            is_first_bible = False
            continue
    
        slide_layout = prs.slides[bible_slide_index].slide_layout
        new_slide = prs.slides.add_slide(slide_layout)
        new_slide_index = bible_slide_index + index

        old_index = prs.slides.index(new_slide)
        xml_slides = prs.slides._sldIdLst  # pylint: disable=W0212
        slides = list(xml_slides)
        xml_slides.remove(slides[old_index])
        xml_slides.insert(new_slide_index, slides[old_index])

        print("title: ", title)
        print("contents: ", contents)

        prs = edit_text_field(
            prs=prs,
            slide_index=new_slide_index,
            is_title=True,
            new_text=title,
        )
        prs = edit_text_field(
            prs=prs,
            slide_index=new_slide_index,
            is_title=False,
            new_text=contents,
            ph_index=10
        )

    return prs

# def remove_slide(prs: PresentationType, slide_index: int) -> PresentationType:
#     del_slide = prs.slides[slide_index]
#     slide_dict = {}

#     for index, value in enumerate(prs.slides._sldIdLst):
#         slide_dict[value.id] = [index, value.rId]

#     slide_id = del_slide.slide_id
#     prs.part.drop_rel(slide_dict[slide_id][1])
#     del prs.slides._sldIdLst[slide_dict[slide_id][0]]

#     return prs


def save_presentation(prs: PresentationType, save_path: str):
    prs.save(save_path)